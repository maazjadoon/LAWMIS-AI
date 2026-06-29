"""
generators/pptx_generator.py
────────────────────────────
Builds a styled LAWMIS KPI PowerPoint presentation from a structured KPI dict.

Design theme:
  • Title slide  — navy background (#1B2A4A), white text, green subtitle bar
  • Content slides — white background, navy title bar, green accent line
  • Data cells    — alternating light-grey rows for readability

Slide layout:
  0 — Title Slide
  1 — Workshops overview
  2 — Licenses overview
  3 — Payments overview
  4 — Emission Tests overview
  5 — RTA Inspections overview
  6 — Top 5 Cities (table)
  7 — Monthly Trends (table)
"""

import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

import config

logger = logging.getLogger(__name__)

# ── Brand colours ─────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)   # #1B2A4A
GREEN   = RGBColor(0x2E, 0xA8, 0x5E)   # #2EA85E
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY_BG = RGBColor(0xF4, 0xF6, 0xF9)   # alternating row bg
GREY_TXT= RGBColor(0x55, 0x65, 0x7A)   # subtle label text

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _solid_fill(shape, colour: RGBColor) -> None:
    """Fill a shape with a solid colour."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_textbox(
    slide, left, top, width, height,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    colour: RGBColor = NAVY,
    align=PP_ALIGN.LEFT,
    word_wrap: bool = True,
) -> None:
    """Add a text box with consistent styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = colour


def _title_bar(slide, title: str) -> None:
    """Add the navy title bar at the top of a content slide."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, Inches(1.0),
    )
    _solid_fill(bar, NAVY)
    bar.line.fill.background()

    # Green accent line below title bar
    line = slide.shapes.add_shape(1, 0, Inches(1.0), SLIDE_W, Inches(0.05))
    _solid_fill(line, GREEN)
    line.line.fill.background()

    _add_textbox(
        slide,
        left=Inches(0.3), top=Inches(0.1),
        width=Inches(12.5), height=Inches(0.8),
        text=title,
        font_size=24, bold=True, colour=WHITE,
    )


def _kpi_block(slide, left, top, label: str, value: str) -> None:
    """Render a single KPI label + value pair."""
    _add_textbox(slide, left, top, Inches(3), Inches(0.35),
                 label, font_size=11, colour=GREY_TXT)
    _add_textbox(slide, left, top + Inches(0.35), Inches(3), Inches(0.5),
                 str(value), font_size=22, bold=True, colour=NAVY)


def _add_table(slide, top, headers: list[str], rows: list[list[str]]) -> None:
    """Add a simple styled table to a slide."""
    from pptx.util import Inches, Pt
    col_count = len(headers)
    col_w = Inches(12.0) / col_count
    table = slide.shapes.add_table(
        len(rows) + 1, col_count,
        Inches(0.65), top,
        Inches(12.0), Inches(0.45 * (len(rows) + 1)),
    ).table

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0] if para.runs else para.add_run()
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE

    # Data rows
    for ri, row in enumerate(rows):
        bg = GREY_BG if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(11)
            run.font.color.rgb = NAVY


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_title(prs: Presentation, date_str: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Full navy background
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    _solid_fill(bg, NAVY)
    bg.line.fill.background()

    # Green accent bar
    bar = slide.shapes.add_shape(1, 0, Inches(3.3), SLIDE_W, Inches(0.08))
    _solid_fill(bar, GREEN)
    bar.line.fill.background()

    _add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(0.9),
                 "LAWMIS KPI Dashboard", 40, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.5),
                 f"Generated: {date_str}", 16, colour=GREEN, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.4),
                 "License & Workshop Management Information System — Analytics Report",
                 13, colour=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)


def _slide_workshops(prs: Presentation, kpi: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "🏭  Workshops Overview")
    w = kpi.get("workshops", {})
    pairs = [
        ("Total Workshops",   w.get("total", "—")),
        ("Active",            w.get("active", "—")),
        ("Pending",           w.get("pending", "—")),
        ("Rejected",          w.get("rejected", "—")),
        ("Approved This Month", w.get("approved_this_month", "—")),
        ("New This Month",    w.get("new_this_month", "—")),
    ]
    cols = 3
    for i, (lbl, val) in enumerate(pairs):
        col = i % cols
        row = i // cols
        _kpi_block(slide,
                   left=Inches(0.5 + col * 4.2),
                   top=Inches(1.4 + row * 1.5),
                   label=lbl, value=val)


def _slide_licenses(prs: Presentation, kpi: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "📄  Licenses Overview")
    lic = kpi.get("licenses", {})
    pairs = [
        ("Total Licenses",        lic.get("total", "—")),
        ("Active",                lic.get("active", "—")),
        ("Expired",               lic.get("expired", "—")),
        ("Expiring in 30 Days",   lic.get("expiring_soon", "—")),
        ("Renewal Count (Total)", lic.get("total_renewals", "—")),
    ]
    for i, (lbl, val) in enumerate(pairs):
        col = i % 3
        row = i // 3
        _kpi_block(slide,
                   left=Inches(0.5 + col * 4.2),
                   top=Inches(1.4 + row * 1.5),
                   label=lbl, value=val)


def _slide_payments(prs: Presentation, kpi: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "💰  Payments & Revenue")
    pay = kpi.get("payments", {})
    pairs = [
        ("Total Revenue",     f"PKR {pay.get('total_revenue', '—')}"),
        ("Revenue This Month",f"PKR {pay.get('revenue_this_month', '—')}"),
        ("Paid Payments",     pay.get("paid_count", "—")),
        ("Pending Payments",  pay.get("pending_count", "—")),
        ("Average Fee",       f"PKR {pay.get('avg_fee', '—')}"),
    ]
    for i, (lbl, val) in enumerate(pairs):
        col = i % 3
        row = i // 3
        _kpi_block(slide,
                   left=Inches(0.5 + col * 4.2),
                   top=Inches(1.4 + row * 1.5),
                   label=lbl, value=val)


def _slide_emission(prs: Presentation, kpi: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "🚗  Emission Tests")
    em = kpi.get("emission_tests", {})
    pairs = [
        ("Tests Today",    em.get("today", "—")),
        ("Total Tests",    em.get("total", "—")),
        ("Pass Rate",      f"{em.get('pass_rate', '—')}%"),
        ("Failure Rate",   f"{em.get('fail_rate', '—')}%"),
        ("Unique Vehicles",em.get("unique_vehicles", "—")),
    ]
    for i, (lbl, val) in enumerate(pairs):
        col = i % 3
        row = i // 3
        _kpi_block(slide,
                   left=Inches(0.5 + col * 4.2),
                   top=Inches(1.4 + row * 1.5),
                   label=lbl, value=val)


def _slide_inspections(prs: Presentation, kpi: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "⭐  RTA Inspections")
    ins = kpi.get("inspections", {})
    pairs = [
        ("Total Inspections", ins.get("total", "—")),
        ("Passed",            ins.get("passed", "—")),
        ("Failed",            ins.get("failed", "—")),
        ("Avg Score",         ins.get("avg_score", "—")),
        ("Follow-up Required",ins.get("follow_up", "—")),
    ]
    for i, (lbl, val) in enumerate(pairs):
        col = i % 3
        row = i // 3
        _kpi_block(slide,
                   left=Inches(0.5 + col * 4.2),
                   top=Inches(1.4 + row * 1.5),
                   label=lbl, value=val)


def _slide_top_cities(prs: Presentation, kpi: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "🏙️  Top Cities by Workshop Count")
    cities = kpi.get("top_cities", [])
    if cities:
        rows = [[c.get("city", ""), str(c.get("count", ""))] for c in cities]
        _add_table(slide, Inches(1.3), ["City", "Workshop Count"], rows)
    else:
        _add_textbox(slide, Inches(1), Inches(2), Inches(10), Inches(1),
                     "No city data available.", 16, colour=GREY_TXT)


def _slide_monthly_trends(prs: Presentation, kpi: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "📈  Monthly Trends (Last 6 Months)")
    trends = kpi.get("monthly_trends", [])
    if trends:
        rows = [
            [
                t.get("month", ""),
                str(t.get("new_workshops", "—")),
                f"PKR {t.get('revenue', '—')}",
                str(t.get("emission_tests", "—")),
            ]
            for t in trends
        ]
        _add_table(slide, Inches(1.3),
                   ["Month", "New Workshops", "Revenue", "Emission Tests"],
                   rows)
    else:
        _add_textbox(slide, Inches(1), Inches(2), Inches(10), Inches(1),
                     "No trend data available.", 16, colour=GREY_TXT)


# ── Public API ────────────────────────────────────────────────────────────────

def generate_pptx(kpi: dict) -> str:
    """
    Generate a styled PPTX from a KPI data dict and save it to EXPORTS_DIR.

    Parameters
    ----------
    kpi : Structured KPI dict produced by handlers/pptx_handler.py.

    Returns
    -------
    Absolute path to the saved .pptx file.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    date_str = datetime.now().strftime("%d %B %Y, %H:%M")

    _slide_title(prs, date_str)
    _slide_workshops(prs, kpi)
    _slide_licenses(prs, kpi)
    _slide_payments(prs, kpi)
    _slide_emission(prs, kpi)
    _slide_inspections(prs, kpi)
    _slide_top_cities(prs, kpi)
    _slide_monthly_trends(prs, kpi)

    # Ensure exports directory exists
    exports_dir = Path(config.EXPORTS_DIR)
    exports_dir.mkdir(parents=True, exist_ok=True)

    filename = f"kpi_dashboard_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    output_path = exports_dir / filename
    prs.save(str(output_path))

    logger.info("PPTX saved: %s", output_path)
    return str(output_path)


def generate_tabular_pptx(headers: list[str], rows: list[list[str]], title: str) -> str:
    """
    Generate a styled PPTX containing slides of tabular custom records.
    Capped at 50 rows (5 slides maximum of tables) to prevent bloat.

    Parameters
    ----------
    headers : List of table column headers
    rows    : List of lists containing cell values (strings)
    title   : Description of the dataset (e.g. 'Workshops in Quetta')

    Returns
    -------
    Absolute path to the saved .pptx file.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    date_str = datetime.now().strftime("%d %B %Y, %H:%M")

    # Slide 0: Cover page
    _slide_title(prs, date_str)

    # Cap rows to 50 for PPTX presentation to keep slides readable
    display_rows = rows[:50]
    chunk_size = 10
    chunks = [display_rows[i:i + chunk_size] for i in range(0, len(display_rows), chunk_size)]

    total_pages = len(chunks)
    for pi, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        page_title = f"📋 {title} (Page {pi + 1} of {total_pages})"
        _title_bar(slide, page_title)

        # Add the chunk data as a table on the slide
        _add_table(slide, Inches(1.3), headers, chunk)

    # Ensure exports directory exists
    exports_dir = Path(config.EXPORTS_DIR)
    exports_dir.mkdir(parents=True, exist_ok=True)

    filename = f"kpi_dashboard_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    output_path = exports_dir / filename
    prs.save(str(output_path))

    logger.info("Tabular PPTX saved: %s", output_path)
    return str(output_path)

