"""
tools/document_tool.py
──────────────────────
Provides text extraction from local PDF and PPTX files.

This is NOT a LangChain @tool — it is a plain utility called by the
document handler. The extracted text is passed directly to the document
agent as context (injected into the system prompt / user message).

Supported formats:
  .pdf  → pypdf
  .pptx → python-pptx
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

MAX_CHARS = 40_000   # limit to avoid overflowing context window


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        parts.append(f"[Page {i + 1}]\n{text.strip()}")

    full_text = "\n\n".join(parts)
    if len(full_text) > MAX_CHARS:
        full_text = full_text[:MAX_CHARS] + "\n\n[... document truncated for context limit ...]"

    logger.info("Extracted %d chars from PDF: %s (%d pages)", len(full_text), file_path, len(reader.pages))
    return full_text


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PPTX file using python-pptx."""
    from pptx import Presentation

    prs = Presentation(file_path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(
                    p.text for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if text:
                    slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_texts))

    full_text = "\n\n".join(parts)
    if len(full_text) > MAX_CHARS:
        full_text = full_text[:MAX_CHARS] + "\n\n[... document truncated for context limit ...]"

    logger.info("Extracted %d chars from PPTX: %s (%d slides)", len(full_text), file_path, len(prs.slides))
    return full_text


def extract_document_text(file_path: str) -> str:
    """
    Auto-detect file type and extract text.

    Parameters
    ----------
    file_path : Absolute or relative path to a PDF or PPTX file.

    Returns
    -------
    Extracted text content as a single string.

    Raises
    ------
    FileNotFoundError   If the file does not exist.
    ValueError          If the file extension is not supported.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(str(path))
    elif ext == ".pptx":
        return extract_text_from_pptx(str(path))
    else:
        raise ValueError(
            f"Unsupported file type '{ext}'. Only .pdf and .pptx are supported."
        )
